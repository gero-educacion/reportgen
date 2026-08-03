import logging
from pptx import Presentation
from pptx.util import Pt
from pathlib import Path
from copy import deepcopy

log = logging.getLogger(__name__)


def determine_template(student: dict, assets_dir: Path) -> Path:
    log.info("determining template...")
    rol = student.get("Rol")

    if rol in ("counseling", "compass-directo"):
        return [
            ("estudiante", assets_dir / "Template_Autoconocimiento 2.0 (counseling).pptx"),
            ("padres",     assets_dir / "Template_Autoconocimiento 2.0 (Padres - Counseling).pptx"),
        ]
    elif rol == "Rojo":
        return [("ccr_rojo",     assets_dir / "CCR_EN BOXES.pptx")]
    elif rol == "Amarillo":
        return [("ccr_amarillo", assets_dir / "CCR_CALENTANDO MOTORES.pptx")]
    elif rol == "Verde":
        return [("ccr_verde",    assets_dir / "CCR_A TODA MARCHA.pptx")]
    elif rol == "autoconfianza":
        return [("autoconfianza", assets_dir / "autoconfianza.pptx")]
    elif rol == "como decidir":
        return [("como decidir",  assets_dir / "como_decidir.pptx")]
    elif rol == "desempatador":
        return [("desempatador",  assets_dir / "desempatador.pptx")]
    elif rol == "contexto":
        return [("contexto", assets_dir / "template_tu_contexto.pptx")]
    elif rol == "UTP":
        return [("estudiante", assets_dir / "template_utp_alumnos.pptx"),
                ("padres", assets_dir / "template_utp_padres.pptx")
            ]
    else:
        return [
            ("estudiante", assets_dir / "Template_Autoconocimiento 2.0 (completo).pptx"),
            ("padres",     assets_dir / "Template_Autoconocimiento 2.0 (PADRES).pptx"),
        ]

def map_placeholders(student: dict) -> dict:
    return {f"<<{k}>>": v for k, v in student.items()}

# ---------------------------------------------------------------------------
# Run-merging helpers
# The core problem: PowerPoint splits a single text run into multiple runs
# whenever formatting changes mid-token, so <<DESC Gest_Temp 02>> might live
# across 3 runs: ["<<DESC", " Gest_Temp", " 02>>"].
# We merge ALL runs in a paragraph into one, do the replacement, then restore
# the original run structure with the first run's formatting.
# ---------------------------------------------------------------------------

def _paragraph_full_text(paragraph) -> str:
    return "".join(r.text for r in paragraph.runs)


def _replace_in_paragraph(paragraph, placeholders: dict) -> bool:
    """
    Merge all runs, replace every placeholder, then rewrite the paragraph.
    Returns True if any replacement was made.
    """
    if not paragraph.runs:
        return False

    full = _paragraph_full_text(paragraph)
    original = full

    for key, val in placeholders.items():
        if key in full:
            if isinstance(val, str):
                full = full.replace(key, val)
            elif isinstance(val, (int, float)):
                full = full.replace(key, f"{round(val * 100, 1)}%")
            # image placeholders are handled separately

    if full == original:
        return False

    # Write the merged text back into the first run, clear the rest
    first_run = paragraph.runs[0]
    first_run.text = full
    for run in paragraph.runs[1:]:
        run.text = ""

    return True


def _find_placeholder_in_paragraph(paragraph, key: str) -> bool:
    return key in _paragraph_full_text(paragraph)


# ---------------------------------------------------------------------------
# Main report generator
# ---------------------------------------------------------------------------

def generate_report(
    student: dict,
    pie_chart_path: Path,
    template_path: Path,
    output_pptx_path: Path,
):
    log.info("Generating PPTX from %s", template_path.name)
    prs = Presentation(template_path)

    placeholders = map_placeholders(student)

    # Separate image placeholders from text/numeric ones
    image_placeholders = {k: v for k, v in placeholders.items()
                          if isinstance(v, dict) and "image" in v}
    text_placeholders  = {k: v for k, v in placeholders.items()
                          if not (isinstance(v, dict) and "image" in v)}

    for slide in prs.slides:
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue

            # ----------------------------------------------------------
            # 1. IMAGE placeholders  (must happen before text replacement
            #    because we remove the shape)
            # ----------------------------------------------------------
            shape_text = shape.text
            for key, val in image_placeholders.items():
                if key not in shape_text:
                    continue

                image_path = Path(val["image"])
                if not image_path.exists():
                    log.warning("Image not found for placeholder %s: %s", key, image_path)
                    continue

                slide.shapes.add_picture(
                    str(image_path),
                    shape.left, shape.top,
                    width=shape.width, height=shape.height,
                )
                # blank the placeholder text so it doesn't show through
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                break   # one image per shape

            # ----------------------------------------------------------
            # 2. Holland chart  (special case: keyed on <<URL GRAFICO HOLLAND>>)
            # ----------------------------------------------------------
            holland_key = "<<URL GRAFICO HOLLAND>>"
            if holland_key in shape.text and pie_chart_path and pie_chart_path.exists():
                slide.shapes.add_picture(
                    str(pie_chart_path),
                    shape.left, shape.top,
                    width=shape.width, height=shape.height,
                )
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                continue

            # ----------------------------------------------------------
            # 3. Text / numeric placeholders  — paragraph-level merge
            # ----------------------------------------------------------
            for paragraph in shape.text_frame.paragraphs:
                replaced = _replace_in_paragraph(paragraph, text_placeholders)
                if replaced:
                    log.debug("Replaced placeholder(s) in paragraph: %r",
                              _paragraph_full_text(paragraph)[:80])

    output_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx_path)
    log.info("PPTX written → %s", output_pptx_path)
    return output_pptx_path