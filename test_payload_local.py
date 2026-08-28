"""
Ad-hoc local test harness — temporary, ported from
.ipynb_checkpoints/test_payload-checkpoint.ipynb into a runnable script.

Mocks every external side effect (Drive, Sheets, SendGrid, SiteGround, MySQL)
and calls process_report_job() directly (same function the RQ worker calls),
with a REAL local Gotenberg for PDF conversion — so this exercises the actual
new concurrent-conversion path in run_student_pipeline.py and the parallel
upload loops in tasks.py, without touching any real external service.
"""
import os
import re
import sys
import time
import random
import logging

ROLE_TO_TEST = os.environ.get("ROLE_TO_TEST", "counseling")

os.environ.setdefault("GOTENBERG_URL", "http://gotenberg:3000")
os.environ.setdefault("SHEETS_STATUS_ID", "fake-sheet-id")
os.environ.setdefault("GOOGLE_CREDENTIALS_JSON", "{}")
os.environ.setdefault("SENDGRID_API_KEY", "fake")
os.environ.setdefault("SENDGRID_FROM", "test@example.com")
os.environ.setdefault("DRIVE_FOLDER_AUTO_EST", "fake-folder-est")
os.environ.setdefault("DRIVE_FOLDER_AUTO_PAD", "fake-folder-pad")

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s.%(msecs)03d | %(message)s",
    datefmt="%H:%M:%S",
)

sys.path.insert(0, os.path.abspath("."))

import app.pipeline.job_config as job_config
import app.tasks as tasks_mod
from pptx import Presentation


def discover_placeholders(template_path):
    prs = Presentation(template_path)
    found = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            found.update(re.findall(r"<<(.*?)>>", shape.text_frame.text))
    return found


HOLLAND_LABELS = ["REALISTA", "INVESTIGATIVO", "ARTÍSTICO", "SOCIAL", "EMPRENDEDOR", "CONVENCIONAL"]
GARDNER_LABELS = [
    "LINGÜÍSTICO-VERBAL", "LÓGICO-MATEMÁTICA", "VISUAL-ESPACIAL", "CORPORAL-KINESTÉSICA",
    "MUSICAL", "INTERPERSONAL", "INTRAPERSONAL", "NATURALISTA",
]
GARDNER_LEVELS = ["ALTA", "MEDIA", "BAJA"]


def _mock_value_for(key):
    return f"Acá iría el texto para el placeholder {key}"


def generate_mock_payload(rol, student_id="test-001", job_id=None, name="OJ Simpson",
                           email="juan@geroeducacion.com", seed=None):
    if seed is not None:
        random.seed(seed)
    job_id = job_id or f"{student_id}-run1"
    job = {
        "student_id": student_id,
        "job_id": job_id,
        "Rol": rol,
        "Nombre y Apellido": name,
        "Email": email,
        "send_email": True,
        "upload_drive": True,
        "post_siteground": True,
        "upload_historic": True,
        "force_rerun": True,
    }

    placeholder_keys = set()
    for suffix, template_path in job_config.get_templates(rol):
        if not template_path.exists():
            print(f"WARNING template not found: {template_path}")
            continue
        placeholder_keys |= discover_placeholders(template_path)
    for key in sorted(placeholder_keys):
        if key not in job:
            job[key] = _mock_value_for(key)

    if job_config.needs_full_pipeline(rol):
        raw_pcts = [random.random() for _ in range(6)]
        total = sum(raw_pcts)
        pcts = [p / total for p in raw_pcts]
        for i, (label, pct) in enumerate(zip(HOLLAND_LABELS, pcts), start=1):
            job[f"Holland {i:02d}"] = label
            job[f"% Holland {i:02d}"] = round(pct, 4)
            job[f"PRE Holland {i:02d}"] = random.choice(GARDNER_LEVELS)
        for i, label in enumerate(GARDNER_LABELS, start=1):
            job[f"Gardner {i:02d}"] = label
            job[f"% Gardner {i:02d}"] = random.choice(GARDNER_LEVELS)
        for i in range(1, 9):
            job[f"% AE {i:02d}"] = round(random.uniform(0.1, 0.95), 2)
        for i in range(1, 5):
            job[f"CARRERA_0{i}"] = f"✔️ Carrera de ejemplo #{i}"

    return job


def mock_upload_pdf_to_drive(pdf_path, target_folder_id, filename):
    print(f"[MOCK Drive] upload  folder={target_folder_id}  file={filename}")
    return f"https://drive.google.com/fake/{filename}"


def mock_upsert_json_to_drive(data, filename, folder_id):
    print(f"[MOCK Drive] upsert-json  folder={folder_id}  file={filename}")
    return f"https://drive.google.com/fake/{filename}"


def mock_send_report_to_siteground(email, drive_link, post_title, description=""):
    print(f"[MOCK SiteGround] notify  email={email}  title={post_title!r}")
    return {"success": True}


def mock_upload_pdf_to_siteground(local_path, filename):
    print(f"[MOCK SiteGround] sftp-upload  file={filename}")
    return f"https://staging2.geroeducacion.com/pdf_storage/{filename}"


def mock_send_report_email(to_email, pdf_paths, student):
    print(f"[MOCK Email] would send to={to_email}  attachments={[str(p) for p in pdf_paths]}")


def mock_update_student_status(student_id, student, status, **kwargs):
    print(f"[MOCK Sheet] row  student_id={student_id}  status={status}  {kwargs}")


def mock_noop(*args, **kwargs):
    print(f"[MOCK] no-op call, args={args} kwargs={kwargs}")


def mock_get_all_links(email):
    return {}


tasks_mod.upload_pdf_to_drive        = mock_upload_pdf_to_drive
tasks_mod.upsert_json_to_drive       = mock_upsert_json_to_drive
tasks_mod.send_report_to_siteground  = mock_send_report_to_siteground
tasks_mod.upload_pdf_to_siteground   = mock_upload_pdf_to_siteground
tasks_mod.send_report_email          = mock_send_report_email
tasks_mod.update_student_status      = mock_update_student_status
tasks_mod.write_majors_to_db         = mock_noop
tasks_mod.post_utp_payload           = mock_noop
tasks_mod.alter_table_reports        = mock_noop
tasks_mod.upsert_historico           = mock_noop
tasks_mod.download_drive_file        = mock_noop
tasks_mod.get_all_links              = mock_get_all_links
tasks_mod.upsert_student             = mock_noop

print("All external side effects mocked — only Gotenberg (real, local) is hit.\n")

ready, problems = job_config.role_is_ready(ROLE_TO_TEST)
print(f"role_is_ready({ROLE_TO_TEST!r}) -> {ready}  problems={problems}")
if not ready:
    sys.exit(1)

job = generate_mock_payload(ROLE_TO_TEST, seed=42)
print(f"Generated payload with {len(job)} fields for role={ROLE_TO_TEST!r}\n")

t0 = time.time()
result = tasks_mod.process_report_job(job)
elapsed = time.time() - t0

print("\n=== RESULT ===")
print(result)
print(f"\nTotal process_report_job() wall time: {elapsed:.1f}s")

job_dir = tasks_mod.APP_TMP_DIR / job["job_id"]
print(f"\njob_dir: {job_dir}")
for f in sorted(job_dir.iterdir()):
    print(f"  {f.name}  ({f.stat().st_size:,} bytes)")
